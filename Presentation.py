import pptx


def read_pptx(file_path):
    """
    Reads a PowerPoint (.pptx) file slide by slide and returns the text content of all slides as a single string.

    :param file_path: The path to the PowerPoint file.
    :type file_path: str
    :return: A string containing the text content of each slide, separated by a divider.
    :rtype: str
    """
    ptx = pptx.Presentation(file_path)
    content = ""
    for slide_number, slide in enumerate(ptx.slides, start=1):

        empty_slide = True

        # Iterate through each shape in the slide
        for shape in slide.shapes:

            # Check if the shape has text
            if hasattr(shape, "text"):

                # add the slide number only if it's not empty
                if empty_slide:
                    empty_slide = False
                    content += f"Slide {slide_number}:\n"

                content += shape.text + '\n'

        if not empty_slide:
            content += "-" * 40 + '\n'

    # remove whitespaces
    stripped_content = content.strip().splitlines()
    return "\n".join(line for line in stripped_content if line.strip())

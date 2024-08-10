import pptx


def read_pptx(file_path):
    """
    Reads a PowerPoint (.pptx) file slide by slide and returns the text content of all slides as a list.

    :param file_path: The path to the PowerPoint file.
    :type file_path: str
    :return: A list containing the text content of each slide.
    :rtype: list<string>
    """
    ptx = pptx.Presentation(file_path)
    slide_content = []
    slides_lst = []
    for slide_number, slide in enumerate(ptx.slides, start=1):
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(f"{shape.text.strip()}\n")

        if slide_content:
            slides_lst.append("\n".join(line for line in slide_content if line.strip()))
            slide_content.clear()

    return slides_lst
